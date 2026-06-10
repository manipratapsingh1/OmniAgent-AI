from typing import Tuple, List, Dict, Optional, Callable
import io
import asyncio
import structlog

from pypdf import PdfReader
from app.rag.chunker import chunk_text, semantic_chunk_text
from app.rag.embeddings import embed_texts
from app.rag.retriever import vector_store, query_vectors

log = structlog.get_logger("ingest")


def _extract_pdf_pypdf(raw: bytes, filename: str) -> str:
    """Extract text from PDF using PyPDF (primary method)"""
    try:
        reader = PdfReader(io.BytesIO(raw))
        text_parts = []
        for page_num, page in enumerate(reader.pages):
            try:
                page_text = page.extract_text()
                if page_text and page_text.strip():
                    text_parts.append(page_text)
            except Exception as page_e:
                log.debug(
                    "pdf.page_extraction_failed_pypdf",
                    filename=filename,
                    page_num=page_num,
                    error=str(page_e)
                )
        
        extracted = "\n".join(text_parts)
        if extracted and extracted.strip():
            log.info(
                "pdf.extracted_pypdf",
                filename=filename,
                text_length=len(extracted),
                pages_extracted=len(text_parts)
            )
            return extracted
        return ""
    except Exception as e:
        log.warning("pdf.pypdf_extraction_failed", filename=filename, error=str(e))
        return ""


def _extract_pdf_pdfplumber(raw: bytes, filename: str) -> str:
    """Extract text from PDF using pdfplumber (fallback - handles compressed PDFs better)"""
    try:
        import pdfplumber
        
        text_parts = []
        with pdfplumber.open(io.BytesIO(raw)) as pdf:
            for page_num, page in enumerate(pdf.pages):
                try:
                    page_text = page.extract_text()
                    if page_text and page_text.strip():
                        text_parts.append(page_text)
                except Exception as page_e:
                    log.debug(
                        "pdf.page_extraction_failed_pdfplumber",
                        filename=filename,
                        page_num=page_num,
                        error=str(page_e)
                    )
        
        extracted = "\n".join(text_parts)
        if extracted and extracted.strip():
            log.info(
                "pdf.extracted_pdfplumber",
                filename=filename,
                text_length=len(extracted),
                pages_extracted=len(text_parts)
            )
            return extracted
        return ""
    except ImportError:
        log.debug("pdfplumber.not_installed")
        return ""
    except Exception as e:
        log.warning("pdf.pdfplumber_extraction_failed", filename=filename, error=str(e))
        return ""


def _extract_docx(raw: bytes, filename: str) -> str:
    """Extract text from DOCX files using python-docx"""
    try:
        import docx
        document = docx.Document(io.BytesIO(raw))
        paragraphs = [p.text for p in document.paragraphs if p.text and p.text.strip()]
        extracted = "\n".join(paragraphs)
        if extracted.strip():
            log.info("docx.extracted", filename=filename, text_length=len(extracted), paragraphs=len(paragraphs))
            return extracted
    except ImportError:
        log.debug("docx.not_installed")
    except Exception as e:
        log.warning("docx.extraction_failed", filename=filename, error=str(e))
    return ""


def _extract_pptx(raw: bytes, filename: str) -> str:
    """Extract text from PPTX slides using python-pptx"""
    try:
        from pptx import Presentation
        prs = Presentation(io.BytesIO(raw))
        text_parts = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text = shape.text
                    if text and text.strip():
                        text_parts.append(text)
        extracted = "\n".join(text_parts)
        if extracted.strip():
            log.info("pptx.extracted", filename=filename, text_length=len(extracted), slides=len(prs.slides))
            return extracted
    except ImportError:
        log.debug("pptx.not_installed")
    except Exception as e:
        log.warning("pptx.extraction_failed", filename=filename, error=str(e))
    return ""


def _extract_xlsx(raw: bytes, filename: str) -> str:
    """Extract text from XLSX files using pandas"""
    try:
        import pandas as pd
        df_dict = pd.read_excel(io.BytesIO(raw), sheet_name=None)
        text_parts = []
        for sheet_name, df in df_dict.items():
            text_parts.append(f"Sheet: {sheet_name}")
            text_parts.append(df.to_string(index=False))
        extracted = "\n\n".join(text_parts)
        if extracted.strip():
            log.info("xlsx.extracted", filename=filename, text_length=len(extracted), sheets=len(df_dict))
            return extracted
    except ImportError:
        log.debug("pandas.not_installed")
    except Exception as e:
        log.warning("xlsx.extraction_failed", filename=filename, error=str(e))
    return ""


def _extract_csv(raw: bytes, filename: str) -> str:
    """Extract text from CSV files using pandas"""
    try:
        import pandas as pd
        df = pd.read_csv(io.BytesIO(raw))
        extracted = df.to_string(index=False)
        if extracted.strip():
            log.info("csv.extracted", filename=filename, text_length=len(extracted))
            return extracted
    except ImportError:
        log.debug("pandas.not_installed")
    except Exception as e:
        log.warning("csv.extraction_failed", filename=filename, error=str(e))
    return ""


def _extract_image_ocr(raw: bytes, filename: str) -> str:
    """Extract text from images using pytesseract (OCR)"""
    try:
        from PIL import Image
        import pytesseract
        img = Image.open(io.BytesIO(raw))
        extracted = pytesseract.image_to_string(img)
        if extracted.strip():
            log.info("image.ocr.extracted", filename=filename, text_length=len(extracted))
            return extracted
    except ImportError:
        log.debug("pytesseract.not_installed")
    except Exception as e:
        log.warning("image.ocr.failed", filename=filename, error=str(e))
    return ""


def _extract_html(raw: bytes, filename: str) -> str:
    """Extract text from HTML files"""
    try:
        from bs4 import BeautifulSoup
        soup = BeautifulSoup(raw, "html.parser")
        extracted = soup.get_text(separator="\n", strip=True)
        if extracted.strip():
            log.info("html.extracted", filename=filename, text_length=len(extracted))
            return extracted
    except ImportError:
        log.debug("bs4.not_installed")
    except Exception as e:
        log.warning("html.extraction_failed", filename=filename, error=str(e))
    return ""


def _extract_text(filename: str, raw: bytes) -> str:
    """Extract text from various file formats with fallbacks"""
    name = filename.lower()

    if name.endswith(".pdf"):
        # Try PyPDF first
        text = _extract_pdf_pypdf(raw, filename)
        if text and text.strip():
            return text

        # Fallback to pdfplumber (better for compressed/complex PDFs)
        log.info("pdf.trying_pdfplumber_fallback", filename=filename)
        text = _extract_pdf_pdfplumber(raw, filename)
        if text and text.strip():
            return text

        log.warning(
            "pdf.no_text_extracted",
            filename=filename,
            method="pypdf and pdfplumber"
        )
        return ""

    if name.endswith(".docx"):
        text = _extract_docx(raw, filename)
        if text and text.strip():
            return text
        log.warning("docx.no_text_extracted", filename=filename)
        return ""

    if name.endswith(".pptx"):
        text = _extract_pptx(raw, filename)
        if text and text.strip():
            return text
        log.warning("pptx.no_text_extracted", filename=filename)
        return ""

    if name.endswith(".xlsx"):
        text = _extract_xlsx(raw, filename)
        if text and text.strip():
            return text
        log.warning("xlsx.no_text_extracted", filename=filename)
        return ""

    if name.endswith(".csv"):
        text = _extract_csv(raw, filename)
        if text and text.strip():
            return text
        log.warning("csv.no_text_extracted", filename=filename)
        return ""

    if name.endswith((".png", ".jpg", ".jpeg", ".webp", ".bmp")):
        text = _extract_image_ocr(raw, filename)
        if text and text.strip():
            return text
        log.warning("image.no_text_extracted", filename=filename)
        return ""

    if name.endswith(('.html', '.htm')):
        text = _extract_html(raw, filename)
        if text and text.strip():
            return text
        log.warning("html.no_text_extracted", filename=filename)
        return ""

    if name.endswith((".md", ".txt", ".markdown")):
        try:
            decoded = raw.decode("utf-8", errors="ignore")
            if decoded and decoded.strip():
                return decoded
            log.warning("text.file_empty", filename=filename)
            return ""
        except Exception as e:
            log.warning("text.decode.failed", filename=filename, error=str(e))
            return ""

    # default: best-effort decode
    try:
        decoded = raw.decode("utf-8", errors="ignore")
        if decoded and decoded.strip():
            return decoded
        log.warning("default.file_empty", filename=filename)
        return ""
    except Exception as e:
        log.warning("default.decode.failed", filename=filename, error=str(e))
        return ""


async def ingest_file(
    user_id: int,
    document_id: int,
    filename: str,
    raw: bytes,
    is_knowledge_base: bool = False,
    progress_callback: Optional[Callable[[int, str], None]] = None,
) -> Tuple[int, int, List[str], List[str]]:
    """
    Ingest a file into the RAG system with batching for large files.
    Returns: (number_of_chunks, number_of_vectors_stored, vector_ids, chunk_texts)
    """
    try:
        log.info(
            "ingest.start",
            document_id=document_id,
            filename=filename,
            file_size_kb=len(raw) / 1024,
        )
        
        text = _extract_text(filename, raw)
        if not text or not text.strip():
            log.warning(
                "ingest.no_text_extracted",
                document_id=document_id,
                filename=filename,
            )
            return 0, 0, [], []

        log.info(
            "ingest.text_extracted",
            document_id=document_id,
            text_length=len(text),
        )
        try:
            if progress_callback:
                progress_callback(10, "text_extracted")
        except Exception:
            pass

        chunks = semantic_chunk_text(text)
        if not chunks:
            log.warning(
                "ingest.no_chunks_created",
                document_id=document_id,
                filename=filename,
            )
            return 0, 0, [], []

        log.info(
            "ingest.chunks_created",
            document_id=document_id,
            chunk_count=len(chunks),
            avg_chunk_size=len(text) // len(chunks) if chunks else 0,
        )
        try:
            if progress_callback:
                progress_callback(30, "chunks_created")
        except Exception:
            pass

        # For large files with many chunks, batch the embeddings to avoid timeouts
        # Increased batch size from 10 to 30 for better throughput
        # Process batches in parallel for 3-5x faster embedding
        max_batch_size = 30  # Embed in larger batches for better performance
        vectors = []
        
        # Create all batch tasks and run in parallel
        batch_tasks = []
        batch_info = []
        
        for batch_start in range(0, len(chunks), max_batch_size):
            batch_end = min(batch_start + max_batch_size, len(chunks))
            batch_chunks = chunks[batch_start:batch_end]
            batch_num = batch_start // max_batch_size + 1
            total_batches = (len(chunks) + max_batch_size - 1) // max_batch_size
            
            batch_info.append({
                "batch_num": batch_num,
                "batch_start": batch_start,
                "batch_end": batch_end,
                "batch_size": len(batch_chunks),
                "total_batches": total_batches,
            })
            
            log.info(
                "ingest.embedding_batch_queued",
                document_id=document_id,
                batch_num=batch_num,
                batch_start=batch_start,
                batch_end=batch_end,
                batch_size=len(batch_chunks),
                total_batches=total_batches,
            )
            
            # Create task for parallel execution (limit to 3 parallel tasks to avoid overwhelming Ollama)
            batch_tasks.append(embed_texts(batch_chunks))
        
        # Execute batches with controlled parallelism (max 3 concurrent requests to Ollama)
        max_parallel = min(3, len(batch_tasks))
        try:
            batch_results = []
            for i in range(0, len(batch_tasks), max_parallel):
                parallel_tasks = batch_tasks[i:i + max_parallel]
                results = await asyncio.gather(*parallel_tasks, return_exceptions=True)
                
                for j, result in enumerate(results):
                    batch_idx = i + j
                    info = batch_info[batch_idx]
                    
                    if isinstance(result, Exception):
                        log.exception(
                            "ingest.batch_embedding_failed",
                            document_id=document_id,
                            batch_range=f"{info['batch_start']}-{info['batch_end']}",
                            error=str(result),
                            error_type=type(result).__name__,
                        )
                        return 0, 0, [], []
                    
                    batch_vectors = result
                    if not batch_vectors or len(batch_vectors) != info['batch_size']:
                        log.error(
                            "ingest.batch_embedding_mismatch",
                            document_id=document_id,
                            batch_range=f"{info['batch_start']}-{info['batch_end']}",
                            expected_vectors=info['batch_size'],
                            actual_vectors=len(batch_vectors) if batch_vectors else 0,
                        )
                        return 0, 0, [], []
                    
                    vectors.extend(batch_vectors)
                    log.info(
                        "ingest.batch_complete",
                        document_id=document_id,
                        batch_num=info['batch_num'],
                        batch_range=f"{info['batch_start']}-{info['batch_end']}",
                        vectors_so_far=len(vectors),
                        total_batches=info['total_batches'],
                    )
                    batch_results.append((batch_idx, batch_vectors))
                    # report progress per batch
                    try:
                        if progress_callback:
                            total = info['total_batches']
                            completed = info['batch_num']
                            # progress between 40 and 90 for embedding batches
                            pct = 40 + int((completed / total) * 50)
                            progress_callback(pct, f"batch_{info['batch_num']}_complete")
                    except Exception:
                        pass
        except Exception as e:
            log.exception(
                "ingest.parallel_embedding_failed",
                document_id=document_id,
                error=str(e),
                error_type=type(e).__name__,
            )
            return 0, 0, [], []

        # Validate all embeddings were created
        if not vectors or len(vectors) != len(chunks):
            log.error(
                "ingest.embedding_mismatch",
                document_id=document_id,
                expected_vectors=len(chunks),
                actual_vectors=len(vectors),
            )
            return 0, 0, [], []
        
        log.info(
            "ingest.embeddings_complete",
            document_id=document_id,
            total_vectors=len(vectors),
        )
        try:
            if progress_callback:
                progress_callback(85, "embeddings_complete")
        except Exception:
            pass

        # Store in vector database
        try:
            metadatas = [
                {
                    "user_id": user_id,
                    "document_id": document_id,
                    "chunk_index": i,
                    "is_knowledge_base": is_knowledge_base,
                }
                for i in range(len(chunks))
            ]
            ids = [f"{document_id}-{i}" for i in range(len(chunks))]
            vector_store.add(
                ids=ids,
                embeddings=vectors,
                documents=chunks,
                metadatas=metadatas,
            )

            # Verify retrieval works: query with one of the embeddings and
            # ensure the stored metadata references this document.
            try:
                verify = query_vectors(embedding=vectors[0], k=1, where={"user_id": user_id}, timeout=5.0)
                verify_docs = (verify.get("documents") or [[]])[0]
                verify_metas = (verify.get("metadatas") or [[]])[0]

                if verify_metas and verify_metas[0].get("document_id") == document_id:
                    log.info(
                        "ingest.success",
                        document_id=document_id,
                        filename=filename,
                        chunks_stored=len(chunks),
                        vectors_stored=len(vectors),
                    )
                    try:
                        if progress_callback:
                            progress_callback(100, "indexed")
                    except Exception:
                        pass
                    return len(chunks), len(vectors), ids, chunks

                # retrieval verification failed
                log.error(
                    "ingest.retrieval_verification_failed",
                    document_id=document_id,
                    verify_metas=verify_metas[:1] if verify_metas else [],
                )
                return 0, 0, [], []
            except Exception as e:
                log.exception(
                    "ingest.retrieval_verification_error",
                    document_id=document_id,
                    error=str(e),
                )
                return 0, 0, [], []
        except Exception as e:
            log.exception(
                "ingest.vector_store_failed",
                document_id=document_id,
                filename=filename,
                error=str(e),
            )
            return 0, 0, [], []

    except Exception as e:
        log.error(
            "ingest.process_failed",
            document_id=document_id,
            filename=filename,
            error=str(e),
        )
        return 0, 0, [], []